from __future__ import annotations
from pathlib import Path
import re
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

EMU = 914400


def _text(sh):
    return getattr(sh, 'text', '') if hasattr(sh, 'text') else ''


def _set_text(sh, value):
    if not hasattr(sh, 'text_frame'):
        return
    value = str(value)
    tf = sh.text_frame
    if not tf.paragraphs:
        tf.text = value
        return
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = value
        for r in p.runs[1:]:
            r.text = ''
    else:
        p.text = value
    for extra in list(tf.paragraphs[1:]):
        el = extra._element
        el.getparent().remove(el)


def _replace(slide, old, new):
    for sh in slide.shapes:
        if _text(sh).strip() == old:
            _set_text(sh, new)
            return True
    return False


def _fill(sh, rgb):
    try:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(rgb)
    except Exception:
        pass


def _font_color(sh, rgb):
    try:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = RGBColor.from_string(rgb)
    except Exception:
        pass


def _set_rag(slide, status):
    # Status is an existing template text shape. Its fill is changed on that same shape.
    for sh in slide.shapes:
        if _text(sh).strip() == '[RAG]':
            _set_text(sh, status)
            _fill(sh, {'GREEN':'00A651','AMBER':'FFC000','RED':'C00000'}.get(status, 'FFC000'))
            _font_color(sh, 'FFFFFF')
            return



def _remove_shapes(shapes):
    """Remove placeholder shapes from the slide without affecting other template geometry."""
    for sh in list(shapes):
        try:
            sp = sh._element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _rgb(hex_color):
    return RGBColor.from_string(hex_color)


def _style_chart(chart, title=None, x_title=None, y_title=None, legend=False, horizontal=False):
    """Apply a restrained, template-faithful chart style with real axes."""
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        try:
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.name = 'Arial'
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb('0058A8')
        except Exception:
            pass
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = True
        try:
            chart.legend.font.name = 'Arial'
            chart.legend.font.size = Pt(7)
        except Exception:
            pass
        try:
            for p in chart.legend.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = 'Arial'; r.font.size = Pt(7)
        except Exception:
            pass
    try:
        chart.chart_style = 10
    except Exception:
        pass
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = _rgb('D9E2F3')
        chart.value_axis.tick_labels.font.name = 'Arial'
        chart.value_axis.tick_labels.font.size = Pt(7)
        chart.value_axis.tick_labels.font.color.rgb = _rgb('404040')
        if y_title:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_title
            chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.name = 'Arial'
            chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
    except Exception:
        pass
    try:
        chart.category_axis.tick_labels.font.name = 'Arial'
        chart.category_axis.tick_labels.font.size = Pt(7)
        chart.category_axis.tick_labels.font.color.rgb = _rgb('404040')
        if x_title:
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.text = x_title
            chart.category_axis.axis_title.text_frame.paragraphs[0].runs[0].font.name = 'Arial'
            chart.category_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(7)
    except Exception:
        pass
    try:
        chart.plot.vary_by_categories = False
    except Exception:
        pass


def _set_series_color(series, color='0058A8'):
    try:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(color)
        series.format.line.color.rgb = _rgb(color)
    except Exception:
        pass


def _add_bar_chart(slide, left, top, width, height, categories, values, *, title='', x_title='', y_title='', horizontal=False, color='0058A8', data_labels=True):
    """Insert a genuine PowerPoint chart in the template's chart box."""
    data = ChartData()
    data.categories = [str(x) for x in categories]
    data.add_series('Cases' if y_title.lower().startswith('case') else 'Count', [float(v or 0) for v in values])
    chart_type = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), data).chart
    _style_chart(chart, title=title, x_title=(y_title if horizontal else x_title), y_title=(x_title if horizontal else y_title), legend=False, horizontal=horizontal)
    _set_series_color(chart.series[0], color)
    if data_labels:
        try:
            chart.series[0].has_data_labels = True
            dl = chart.series[0].data_labels
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
            dl.font.name = 'Arial'; dl.font.size = Pt(7); dl.font.color.rgb = _rgb('404040')
            dl.show_value = True
        except Exception:
            pass
    try:
        chart.value_axis.minimum_scale = 0
        mx = max([float(v or 0) for v in values] + [1])
        chart.value_axis.maximum_scale = max(mx * 1.15, 1)
        chart.value_axis.major_unit = max(1, round(chart.value_axis.maximum_scale / 5))
    except Exception:
        pass
    try:
        chart.category_axis.reverse_order = True if horizontal else False
    except Exception:
        pass
    return chart


def _add_stacked_ward_chart(slide, left, top, width, height, areas, covered, missing):
    data = ChartData()
    data.categories = [str(x) for x in areas]
    data.add_series('Covered', [float(x or 0) for x in covered])
    data.add_series('Missing', [float(x or 0) for x in missing])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(left), Inches(top), Inches(width), Inches(height), data).chart
    _style_chart(chart, title='', x_title='Area', y_title='Wards', legend=True, horizontal=True)
    _set_series_color(chart.series[0], '0058A8')
    _set_series_color(chart.series[1], 'FFC000')
    for ser in chart.series:
        try:
            ser.has_data_labels = True
            ser.data_labels.show_value = True
            ser.data_labels.font.name = 'Arial'; ser.data_labels.font.size = Pt(7); ser.data_labels.font.color.rgb = _rgb('404040')
        except Exception:
            pass
    try:
        chart.value_axis.minimum_scale = 0
        mx = max([float(a or 0)+float(b or 0) for a,b in zip(covered,missing)] + [1])
        chart.value_axis.maximum_scale = max(mx * 1.12, 1)
        chart.value_axis.major_unit = max(1, round(chart.value_axis.maximum_scale / 5))
    except Exception:
        pass
    return chart


def _add_hourly_chart(slide, left, top, width, height, hours, values):
    labels = [f'{int(h):02d}:00' for h in hours]
    return _add_bar_chart(slide, left, top, width, height, labels, values, title='', x_title='Hour of day', y_title='Cases', horizontal=False, color='0058A8')



def _add_line_chart(slide, left, top, width, height, hours, values):
    """Insert a genuine PowerPoint line chart for cases logged by hour."""
    data = ChartData()
    labels = [f'{int(h):02d}:00' for h in hours]
    data.categories = labels
    data.add_series('Cases', [float(v or 0) for v in values])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left), Inches(top), Inches(width), Inches(height), data
    ).chart
    _style_chart(chart, title='', x_title='Hour of day', y_title='Cases', legend=False, horizontal=False)
    _set_series_color(chart.series[0], '0058A8')
    try:
        chart.series[0].marker.style = 2
        chart.series[0].marker.size = 6
    except Exception:
        pass
    try:
        chart.value_axis.minimum_scale = 0
        mx = max([float(v or 0) for v in values] + [1])
        chart.value_axis.maximum_scale = max(mx * 1.20, 1)
        chart.value_axis.major_unit = max(1, round(chart.value_axis.maximum_scale / 5))
    except Exception:
        pass
    return chart

def _set_by_index(slide, indices, values):
    for i, v in zip(indices, values):
        if i < len(slide.shapes):
            _set_text(slide.shapes[i], v)


def _cleanup_instructions(prs):
    instruction_phrases = [
        'Use this as the single front page for all corridors. Keep only the six headline measures relevant to management.',
        'Replace with cases and new wards by day.',
        'Replace with stacked bars by municipality/region.',
        'Replace with cases by channel or channel share.',
        'Replace with cases by hour and cumulative new wards.',
        'Replace with top 5-8 services by cases.',
        'Replace with missing wards, VOC queue and key data exceptions.',
        'Template: replace all [bracketed] fields and sample chart data',
    ]
    for slide in prs.slides:
        for sh in slide.shapes:
            t = _text(sh).strip()
            if t in instruction_phrases:
                _set_text(sh, '')


def _generic_placeholder_cleanup(prs):
    exact = {
        '[##/target]':'0/0', '[##/##]':'0/0', '[##%]':'0%', '[##]':'0', '[target]':'0',
        '[summary]':'0', '[key note]':'0', '[Status]':'Not supplied', '[RAG]':'AMBER',
        '[Case ref]':'0', '[Ward]':'0', '[Service]':'Not available', '[Call / verify / close]':'Call / verify / close',
        '[Ops lead]':'Ops lead', '[Next close]':'Next close', '[Supervisor]':'Supervisor', '[Today]':'Today',
        '[hours covered]':'Observed window', '[VOC lead]':'VOC lead', '[This week]':'Current period',
        '[##] new wards':'0 new wards', '[exceptions corrected]':'0 exceptions corrected',
    }
    for slide in prs.slides:
        for sh in slide.shapes:
            if not hasattr(sh, 'text_frame'): continue
            t = _text(sh)
            new = t
            for a,b in exact.items(): new = new.replace(a,b)
            new = re.sub(r'\[##/target\]', '0/0', new)
            new = re.sub(r'\[##/##\]', '0/0', new)
            new = re.sub(r'\[##%\]', '0%', new)
            new = re.sub(r'\[##\]', '0', new)
            if new != t: _set_text(sh, new)


def _fill_channel_table(slide, ch):
    rows = ch.head(4).reset_index(drop=True) if isinstance(ch, pd.DataFrame) else pd.DataFrame()
    names = [24,29,34,39]
    cases = [25,30,35,40]
    shares = [26,31,36,41]
    new_wards = [27,32,37,42]
    notes = [28,33,38,43]
    for i in range(4):
        if i < len(rows):
            r=rows.iloc[i]
            _set_text(slide.shapes[names[i]], str(r['value']))
            _set_text(slide.shapes[cases[i]], f"{int(r['cases']):,}")
            _set_text(slide.shapes[shares[i]], f"{float(r['share_pct']):.1f}%")
            _set_text(slide.shapes[new_wards[i]], f"{int(r.get('new_wards',0) or 0):,}")
            _set_text(slide.shapes[notes[i]], 'Demand concentration' if i==0 else 'Monitor mix')
        else:
            _set_text(slide.shapes[names[i]], '0')
            _set_text(slide.shapes[cases[i]], '0')
            _set_text(slide.shapes[shares[i]], '0%')
            _set_text(slide.shapes[new_wards[i]], '0')
            _set_text(slide.shapes[notes[i]], '0')

def _fill_service_table(slide, cat):
    rows = cat.head(3).reset_index(drop=True) if isinstance(cat, pd.DataFrame) else pd.DataFrame()
    names=[31,35,39,43]; cases=[32,36,40,44]; shares=[33,37,41,45]; notes=[34,38,42,46]
    total = int(cat['cases'].sum()) if isinstance(cat,pd.DataFrame) and not cat.empty else 0
    for i in range(4):
        if i < len(rows):
            r=rows.iloc[i]; name=r['value']; c=int(r['cases']); pct=float(r['share_pct'])
            note='Leading demand category' if i==0 else 'Demand indicator'
        elif i == 3 and total:
            shown=int(rows['cases'].sum()) if len(rows) else 0
            c=max(total-shown,0); name='Other'; pct=(c/total*100 if total else 0); note='Remaining categories'
        else:
            name='0'; c=0; pct=0; note='0'
        _set_text(slide.shapes[names[i]], name); _set_text(slide.shapes[cases[i]], f'{c:,}'); _set_text(slide.shapes[shares[i]], f'{pct:.1f}%'); _set_text(slide.shapes[notes[i]], note)

def generate_powerpoint(template_path, output_path, analysis, intelligence, municipality, reporting_date, period_covered, close_time, audience_view=None, voc_analysis=None):
    prs=Presentation(str(template_path))
    ssum=analysis['summary']; corridor=municipality or 'OOP Corridor'
    date_str=reporting_date.strftime('%d %B %Y') if hasattr(reporting_date,'strftime') else str(reporting_date)
    close=str(close_time)
    total=int(ssum.get('total_wards') or 0); covered=int(ssum.get('distinct_wards') or 0); missing=int(ssum.get('missing_wards') or 0)
    cov=float(ssum.get('ward_coverage_pct') or 0)
    voc_n=int(voc_analysis.get('responses',0)) if voc_analysis else 0
    d=analysis.get('dates',{})
    sd=d.get('selected_day',{}) if isinstance(d,dict) else {}
    daily_cases=int(sd.get('cases',0) or 0)
    daily_new=int(sd.get('new_wards',0) or 0)
    running_wards=int(sd.get('running_wards',covered) or 0)
    still_needed=int(sd.get('still_needed',missing) or 0)
    daily_coverage=float(sd.get('coverage_pct',cov) or 0)
    status='GREEN' if ssum.get('invalid_case_records',0)==0 and (ssum.get('ward_coverage_pct') in (None,100)) else 'AMBER'

    # Common template content only.
    for slide in prs.slides:
        for sh in slide.shapes:
            if not hasattr(sh,'text_frame'): continue
            t=_text(sh)
            if '[CORRIDOR NAME]' in t or '[MUNICIPALITIES / REGIONS]' in t:
                _set_text(sh,t.replace('[CORRIDOR NAME]',corridor).replace('[MUNICIPALITIES / REGIONS]',corridor))
            if '[REPORTING DATE]' in _text(sh): _set_text(sh,_text(sh).replace('[REPORTING DATE]',date_str))
            if '[PERIOD COVERED]' in _text(sh): _set_text(sh,_text(sh).replace('[PERIOD COVERED]',period_covered or 'Not specified'))
            if '[DATA CLOSE TIME]' in _text(sh): _set_text(sh,_text(sh).replace('[DATA CLOSE TIME]',close))
            if 'INTERNAL SDI OPERATIONS REPORTING' in _text(sh):
                _set_text(sh,f'{corridor} | INTERNAL SDI OPERATIONS REPORTING | {date_str} | {close}')
        _set_rag(slide,status)

    # Slide 1
    s=prs.slides[0]
    _set_by_index(s,[6,9,12,15,18,21],[f"{daily_cases:,}",f'{running_wards}/{total}',f'{daily_coverage:.1f}%',str(still_needed),str(voc_n) if voc_analysis else '0',status])
    # Whole status card is the template's existing status shape.
    _set_rag(s,status)

    # Slide 2
    s=prs.slides[1]
    _set_text(s.shapes[11],f"{daily_cases:,} valid cases on {date_str}; {daily_new} new wards today; {running_wards} wards reached cumulatively ({daily_coverage:.1f}% of {total}).")
    _set_text(s.shapes[14],f"Today: {daily_cases:,} cases and {daily_new} new wards. New wards are unique to the selected reporting date after checking all prior valid dates; running wards are cumulative unique wards.")
    pace=int(sd.get('suggested_new_wards_per_remaining_day',0) or 0)
    rem_days=int(sd.get('remaining_workdays_in_week',0) or 0)
    if still_needed > 0 and rem_days > 0:
        attention=f"{still_needed} ward allocations remain. To reach the 112-ward weekly target by Friday, aim for {pace} new wards per remaining working day ({rem_days} days)."
    elif still_needed > 0:
        attention=f"{still_needed} ward allocations remain against the supplied 112-ward weekly target."
    else:
        attention="The 112-ward target has been reached cumulatively through the selected reporting date."
    _set_text(s.shapes[17],attention)
    _set_text(s.shapes[19],'Cases logged per hour')
    day=d.get('by_date',pd.DataFrame()) if d.get('available') else pd.DataFrame()
    day=day.reset_index(drop=True) if isinstance(day,pd.DataFrame) else pd.DataFrame()
    tracker=d.get('daily_tracker',pd.DataFrame()) if d.get('available') else pd.DataFrame()
    tracker=tracker.tail(4).reset_index(drop=True) if isinstance(tracker,pd.DataFrame) else pd.DataFrame()
    for i,row in enumerate([42,47,52,57]):
        if i<len(tracker):
            rr=tracker.iloc[i]
            _set_text(s.shapes[row],str(rr['date'])); _set_text(s.shapes[row+1],f"{int(rr['cases']):,}")
            _set_text(s.shapes[row+2],f"{int(rr['new_wards']):,}"); _set_text(s.shapes[row+3],f"{int(rr['running_wards']):,}/{total}"); _set_text(s.shapes[row+4],str(int(rr['still_needed'])) if pd.notna(rr['still_needed']) else 'Not configured')
        else:
            for j in range(5): _set_text(s.shapes[row+j],'—')
    _set_text(s.shapes[63],f"Validate ward representation ({missing} missing).\nProtect observed peak intake periods.\nUse the addendum for detailed evidence.")
    _remove_shapes([s.shapes[i] for i in [20,21,22,23,24] if i < len(s.shapes)])
    hourly=d.get('by_hour',pd.DataFrame()) if d.get('available') else pd.DataFrame()
    hourly=hourly.reset_index(drop=True) if isinstance(hourly,pd.DataFrame) else pd.DataFrame()
    if hourly.empty:
        hours=[0]
        hourly_values=[0]
    else:
        hours=hourly['hour'].tolist()
        hourly_values=hourly['cases'].tolist()
    _add_line_chart(s,0.75,3.88,3.75,1.55,hours,hourly_values)

    # Slide 3
    s=prs.slides[2]
    cc=analysis.get('corridor_coverage',pd.DataFrame())
    rows=[]
    if isinstance(cc,pd.DataFrame) and not cc.empty:
        for _,r in cc.head(2).iterrows():
            rows.append((str(r['corridor']),int(r['target']),int(r['covered']),int(r['missing']),f"{float(r['coverage_pct']):.1f}%"))
    while len(rows)<2: rows.append(('—',0,0,0,'0%'))
    daily_corridor_covered = int(cc['covered'].sum()) if isinstance(cc,pd.DataFrame) and not cc.empty else int(sd.get('daily_wards',0) or 0)
    daily_corridor_target = int(cc['target'].sum()) if isinstance(cc,pd.DataFrame) and not cc.empty else total
    daily_corridor_missing = max(daily_corridor_target - daily_corridor_covered, 0)
    daily_corridor_pct = (daily_corridor_covered / daily_corridor_target * 100) if daily_corridor_target else 0
    rows.append(('Total',daily_corridor_target,daily_corridor_covered,daily_corridor_missing,f'{daily_corridor_pct:.1f}%'))
    starts=[24,29,34]
    for base,row in zip(starts,rows):
        for idx,val in zip([base,base+1,base+2,base+3,base+4],row): _set_text(s.shapes[idx],val)
    for base in [39]:
        unique_master=int(analysis.get('ward_master',{}).get('unique_wards',total) or total)
        outside=int(analysis.get('summary',{}).get('ward_master_outside_unique_wards',0) or 0)
        for idx,val in zip([base,base+1,base+2,base+3,base+4],('Validation',total,unique_master,outside,f"{unique_master} unique / {total} listed")): _set_text(s.shapes[idx],val)
    _set_text(s.shapes[45],'Covered vs missing wards by area')
    _set_text(s.shapes[55],f"Selected-day coverage: {daily_corridor_covered} of {daily_corridor_target} listed ward allocations were covered across North East and South East. Cumulative progress remains {running_wards}/{total}; master validation: {analysis.get('summary',{}).get('ward_master_unique_wards',total)} unique ward numbers and {analysis.get('summary',{}).get('ward_master_outside_unique_wards',0)} observed ward numbers outside the master.")
    _set_text(s.shapes[58],'The addendum contains the ward evidence and validation trail.')
    _remove_shapes([s.shapes[i] for i in [46,47,48,49,50] if i < len(s.shapes)])
    _add_stacked_ward_chart(s,6.72,2.45,5.65,1.95,cc['corridor'].head(2).tolist() if isinstance(cc,pd.DataFrame) and not cc.empty else [corridor],[int(x) for x in cc['covered'].head(2)] if isinstance(cc,pd.DataFrame) and not cc.empty else [running_wards],[int(x) for x in cc['missing'].head(2)] if isinstance(cc,pd.DataFrame) and not cc.empty else [still_needed])

    # Slide 4
    s=prs.slides[3]; ch=analysis.get('channel'); _fill_channel_table(s,ch)
    ch_chart = ch.head(5).reset_index(drop=True) if isinstance(ch,pd.DataFrame) else pd.DataFrame()
    _set_text(s.shapes[45],'Channel mix')
    hourly=d.get('by_hour',pd.DataFrame()) if d.get('available') else pd.DataFrame()
    h=hourly.sort_values('hour').reset_index(drop=True) if isinstance(hourly,pd.DataFrame) else pd.DataFrame()
    first=d.get('day_min_time'); last=d.get('day_max_time'); peak=d.get('peak_hour'); peakn=d.get('peak_hour_cases',0)
    if first is not None and last is not None:
        note=f'Observed window: {first:%H:%M}–{last:%H:%M}. Peak: {int(peak):02d}:00 ({int(peakn):,} cases).'
    else: note='Observed hourly activity: 0 cases.'
    _set_text(s.shapes[54],'Time coverage / hourly activity')
    _set_text(s.shapes[64],note)
    _remove_shapes([s.shapes[i] for i in [46,47,48,49,50,55,56,57,58,59] if i < len(s.shapes)])
    _add_bar_chart(s,7.05,2.40,5.45,1.35,ch_chart['value'].tolist() if not ch_chart.empty else ['No data'],ch_chart['cases'].tolist() if not ch_chart.empty else [0],x_title='Channel',y_title='Cases',horizontal=False,color='0058A8')
    _add_hourly_chart(s,0.70,4.62,5.75,1.18,h['hour'].tolist() if not h.empty else [0],h['cases'].tolist() if not h.empty else [0])

    # Slide 5
    s=prs.slides[4]; cat=analysis.get('category1'); _fill_service_table(s,cat)
    cat_chart = cat.head(5).reset_index(drop=True) if isinstance(cat,pd.DataFrame) else pd.DataFrame()
    _set_text(s.shapes[10],'Top services logged')
    _set_text(s.shapes[49],'Use demand concentration to target follow-up; do not interpret case volume as departmental performance.')
    _set_text(s.shapes[52],'Use the addendum for low-volume services and ward-level detail.')
    _remove_shapes([s.shapes[i] for i in [11,12,13,14,15] if i < len(s.shapes)])
    _add_bar_chart(s,0.70,2.48,5.00,1.70,cat_chart['value'].tolist() if not cat_chart.empty else ['No data'],cat_chart['cases'].tolist() if not cat_chart.empty else [0],x_title='Cases',y_title='Service category',horizontal=True,color='0058A8')

    # Slide 6
    s=prs.slides[5]
    if voc_analysis:
        vals=[voc_analysis.get('responses',0),voc_analysis.get('happiness','0%'),voc_analysis.get('resolved',0),voc_analysis.get('survey_ready',0),'Available']
    else: vals=[0,'0%',0,'0/0','Not supplied']
    _set_by_index(s,[10,13,16,19,22],vals)
    if not voc_analysis:
        _set_text(s.shapes[50],'VOC source not supplied. Reported VOC metrics are 0 or 0/0 and no citizen sentiment is inferred.')
    else:
        _set_text(s.shapes[50],'VOC source supplied. Metrics are calculated only from supported response fields; detailed responses remain in the addendum.')

    # Slide 7
    s=prs.slides[6]
    action=[('Missing wards','Ops lead','Next close',f'{missing} wards validated/assigned'),('Channel/time gap','Supervisor','Next operating window','Peak hours protected'),('VOC queue','VOC lead','Current period',f'{voc_n} VOC responses reviewed' if voc_analysis else '0 VOC responses'),('Data quality','Data lead','Before circulation',f"{len(analysis.get('quality',[]))} quality findings reviewed")]
    for base,row in zip([22,26,30,34],action):
        for idx,val in zip([base,base+1,base+2,base+3],row): _set_text(s.shapes[idx],val)
    control_categories=['Missing wards','Channel/time','VOC queue','Data quality']
    control_values=[missing,1 if isinstance(ch,pd.DataFrame) and not ch.empty else 0,voc_n,len(analysis.get('quality',[]))]
    _set_text(s.shapes[48],'Decision needed')
    _set_text(s.shapes[49],f'Confirm and assign the {missing} unrepresented wards before the next operational close.')
    _set_text(s.shapes[39],'Priority control load')
    _set_text(s.shapes[51],'Data caveat')
    _set_text(s.shapes[52],'Area-level municipality breakdown requires a validated ward/master reference mapping.')
    _remove_shapes([s.shapes[i] for i in [40,41,42,43,44] if i < len(s.shapes)])
    _add_bar_chart(s,7.25,2.45,5.05,2.15,control_categories,control_values,x_title='Count',y_title='Control area',horizontal=True,color='0058A8')

    _cleanup_instructions(prs)
    _generic_placeholder_cleanup(prs)
    # Any remaining obvious instruction/footer placeholders are removed rather than shown.
    for slide in prs.slides:
        for sh in slide.shapes:
            if not hasattr(sh,'text_frame'): continue
            t=_text(sh)
            if '[' in t and ']' in t:
                # Never expose template tokens in the final report.
                if re.search(r'\[(?:##|target|summary|Service|Municipality|Case|Ward|Status|RAG|Today|Final|Ops|Supervisor|VOC|Data|Call)', t):
                    _set_text(sh,'0')
    Path(output_path).parent.mkdir(parents=True,exist_ok=True)
    prs.save(str(output_path))
    return Path(output_path)
