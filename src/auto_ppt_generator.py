from __future__ import annotations
from pathlib import Path
from typing import Any
import math
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import ChartData

NAVY='17324D'; BLUE='0058A8'; GOLD='FFC000'; LIGHT='F3F6FA'; MID='D9E2F3'; DARK='263238'; WHITE='FFFFFF'; RED='C00000'; GREEN='2E7D32'
FONT='Aptos'


def rgb(x): return RGBColor.from_string(x)

def box(slide, x,y,w,h, fill=WHITE, line=MID, radius=True):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=rgb(fill)
    shp.line.color.rgb=rgb(line); shp.line.width=Pt(0.7)
    return shp

def text(slide, x,y,w,h, value, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=tb.text_frame; tf.clear(); tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(value); r.font.name=FONT; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(color)
    return tb

def title_bar(slide, title, subtitle=''):
    text(slide,0.55,0.32,11.7,0.45,title,23,True,NAVY)
    if subtitle: text(slide,0.57,0.80,11.4,0.3,subtitle,9,False,'5C6B78')
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.16), Inches(11.75), Inches(0.04)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb=rgb(GOLD); slide.shapes[-1].line.fill.background()

def footer(slide, dept, team, period):
    text(slide,0.55,7.10,8.0,0.22,f"{dept}  •  {team}  •  {period}",7,False,'6B7785')
    text(slide,10.6,7.10,1.7,0.22,"Evidence-led report",7,False,'6B7785',PP_ALIGN.RIGHT)

def add_chart(slide, data, x,y,w,h, title, horizontal=True):
    if not data: return
    cats=[str(r.get('value',''))[:28] for r in data]
    vals=[float(r.get('cases',0) or 0) for r in data]
    cd=ChartData(); cd.categories=cats; cd.add_series('Cases',vals)
    typ=XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart=slide.shapes.add_chart(typ, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_title=True; chart.chart_title.text_frame.text=title
    chart.chart_title.text_frame.paragraphs[0].runs[0].font.name=FONT; chart.chart_title.text_frame.paragraphs[0].runs[0].font.size=Pt(10); chart.chart_title.text_frame.paragraphs[0].runs[0].font.bold=True; chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb=rgb(NAVY)
    chart.has_legend=False
    ser=chart.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb=rgb(BLUE); ser.format.line.color.rgb=rgb(BLUE)
    try:
        ser.has_data_labels=True; ser.data_labels.show_value=True; ser.data_labels.font.name=FONT; ser.data_labels.font.size=Pt(7); ser.data_labels.font.color.rgb=rgb(DARK); ser.data_labels.position=XL_LABEL_POSITION.OUTSIDE_END
    except Exception: pass
    try:
        chart.value_axis.minimum_scale=0; chart.value_axis.maximum_scale=max(max(vals+[1])*1.18,1); chart.value_axis.major_gridlines.format.line.color.rgb=rgb(MID)
        chart.value_axis.tick_labels.font.name=FONT; chart.value_axis.tick_labels.font.size=Pt(7)
        chart.category_axis.tick_labels.font.name=FONT; chart.category_axis.tick_labels.font.size=Pt(7)
        if horizontal: chart.category_axis.reverse_order=True
    except Exception: pass
    return chart

def add_line(slide, rows, x,y,w,h):
    if not rows: return
    rows=sorted(rows,key=lambda r:r['date']); cats=[pd.Timestamp(r['date']).strftime('%d %b') for r in rows]; vals=[float(r['cases']) for r in rows]
    cd=ChartData(); cd.categories=cats; cd.add_series('Cases',vals)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    chart.has_title=False; chart.has_legend=False
    ser=chart.series[0]; ser.format.line.color.rgb=rgb(BLUE); ser.format.line.width=Pt(2)
    try: ser.marker.size=6
    except Exception: pass
    try:
        chart.value_axis.minimum_scale=0; chart.value_axis.major_gridlines.format.line.color.rgb=rgb(MID); chart.value_axis.tick_labels.font.name=FONT; chart.value_axis.tick_labels.font.size=Pt(7); chart.category_axis.tick_labels.font.name=FONT; chart.category_axis.tick_labels.font.size=Pt(7)
    except Exception: pass
    return chart

def add_kpi(slide,x,y,w,label,value,accent=BLUE):
    box(slide,x,y,w,1.05,WHITE,MID)
    text(slide,x+0.16,y+0.14,w-0.3,0.25,label,8,False,'687783')
    text(slide,x+0.16,y+0.40,w-0.3,0.43,value,20,True,accent)

def _finding_text(row):
    return str(row.get('finding',''))

def add_donut(slide, data, x, y, w, h):
    if not data: return
    cats=[str(r.get('value',''))[:22] for r in data]
    vals=[float(r.get('cases',0) or 0) for r in data]
    cd=ChartData(); cd.categories=cats; cd.add_series('Cases', vals)
    ch=slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.RIGHT
    try:
        ch.legend.font.name=FONT; ch.legend.font.size=Pt(8)
    except Exception: pass
    try:
        ch.plots[0].has_data_labels=True; ch.plots[0].data_labels.show_percentage=True
        ch.plots[0].data_labels.font.name=FONT; ch.plots[0].data_labels.font.size=Pt(8)
    except Exception: pass
    return ch



def generate_auto_powerpoint(out_path: Path, blueprint: dict[str,Any]):
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    dept=blueprint.get('department','Report'); team=blueprint.get('requesting_team',''); period=blueprint.get('period_label','Selected period'); audience=blueprint.get('audience','')
    evidence=blueprint.get('period_evidence',{}) or {}
    for s in blueprint.get('slides',[]):
        slide=prs.slides.add_slide(blank); typ=s['type']
        if typ=='cover':
            bg=slide.background.fill; bg.solid(); bg.fore_color.rgb=rgb(NAVY)
            text(slide,0.75,1.25,11.5,0.65,s['title'],30,True,WHITE)
            text(slide,0.78,2.05,10.8,0.55,s['subtitle'],18,False,'DCE8F2')
            text(slide,0.78,3.0,7.5,0.35,f"Requesting team: {team}",11,False,'DCE8F2')
            text(slide,0.78,3.45,7.5,0.35,f"Audience: {audience}",11,False,'DCE8F2')
            text(slide,0.78,5.95,10.5,0.3,"Automatically designed from the selected data, period and reporting brief",9,False,'B9C8D6')
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(5.62), Inches(1.4), Inches(0.05)).fill.solid(); slide.shapes[-1].fill.fore_color.rgb=rgb(GOLD); slide.shapes[-1].line.fill.background()
            continue
        title_bar(slide,s.get('title','Report'), s.get('subtitle',''))
        if typ=='executive_summary':
            rows=s.get('findings',[])
            add_kpi(slide,0.6,1.45,2.25,'Period cases',f"{evidence.get('case_count',0):,}")
            add_kpi(slide,3.0,1.45,2.25,'Distinct wards',str(evidence.get('distinct_wards','Not available')))
            add_kpi(slide,5.4,1.45,2.25,'Period',blueprint.get('period_type',''))
            comp=(s.get('comparison') or {})
            change=comp.get('change_pct')
            add_kpi(slide,7.8,1.45,2.25,'Vs comparison',f"{change:+.1f}%" if isinstance(change,(int,float)) else 'Not available', RED if isinstance(change,(int,float)) and change<0 else BLUE)
            box(slide,0.6,2.8,12.1,3.65,LIGHT,LIGHT)
            text(slide,0.82,3.02,11.5,0.3,'What matters',13,True,NAVY)
            y=3.48
            for r in rows[:5]:
                text(slide,0.85,y,0.45,0.32,str(r.get('id','•')),9,True,GOLD)
                text(slide,1.28,y,10.95,0.55,_finding_text(r),10,False,DARK); y+=0.58
        elif typ=='kpi_snapshot':
            metrics=s.get('metrics',[])
            widths=11.4/max(len(metrics),1)
            for i,(label,value) in enumerate(metrics):
                add_kpi(slide,0.75+i*widths,1.7,widths-0.18,label,value,BLUE if i==0 else NAVY)
            box(slide,0.75,3.15,11.55,2.65,LIGHT,LIGHT)
            text(slide,1.0,3.48,10.9,0.3,'Design decision',13,True,NAVY)
            text(slide,1.0,3.95,10.8,1.25,'The automatic designer selected a KPI snapshot because the selected period does not contain enough daily observations for a meaningful trend chart.',10,False,DARK)
            text(slide,1.0,5.35,10.8,0.3,'No trend is inferred from a single observation.',9,False,'5C6B78')
        elif typ=='trend':
            add_line(slide,s.get('data',[]),0.7,1.55,8.15,4.75)
            comp=s.get('comparison') or {}
            box(slide,9.15,1.55,3.45,4.75,LIGHT,LIGHT)
            text(slide,9.4,1.85,2.9,0.3,'Period position',12,True,NAVY)
            text(slide,9.4,2.4,2.7,0.25,'Cases',8,False,'687783'); text(slide,9.4,2.68,2.7,0.42,f"{evidence.get('case_count',0):,}",22,True,BLUE)
            if comp.get('comparison_case_count') is not None:
                text(slide,9.4,3.35,2.7,0.25,'Comparison cases',8,False,'687783'); text(slide,9.4,3.63,2.7,0.42,f"{comp['comparison_case_count']:,}",18,True,NAVY)
            text(slide,9.4,4.35,2.7,0.25,'Interpretation',8,False,'687783')
            note='Use the trend to identify concentration or change; do not infer causality from volume alone.'
            text(slide,9.4,4.62,2.75,1.0,note,9,False,DARK)
        elif typ=='breakdown':
            if s.get('chart_type')=='donut':
                add_donut(slide,s.get('data',[]),0.8,1.5,11.3,4.8)
            else:
                add_chart(slide,s.get('data',[]),0.7,1.55,11.9,4.85,s.get('subtitle','Breakdown'))
        elif typ=='coverage':
            val=s.get('value'); total=s.get('total')
            display=f"{val:,}" if isinstance(val,int) else 'Not available'
            add_kpi(slide,0.75,1.65,3.0,'Wards represented',display,BLUE)
            if total:
                pct=(float(val)/float(total)*100) if isinstance(val,(int,float)) else None
                add_kpi(slide,3.95,1.65,3.0,'Configured total',f"{total:,}",NAVY)
                add_kpi(slide,7.15,1.65,3.0,'Coverage',f"{pct:.1f}%" if pct is not None else 'Not available',GREEN if pct==100 else BLUE)
            rows=s.get('corridor',[])
            if rows:
                cats=[r.get('corridor','') for r in rows]; covered=[float(r.get('covered',0)) for r in rows]; missing=[float(r.get('missing',0)) for r in rows]
                cd=ChartData(); cd.categories=cats; cd.add_series('Covered',covered); cd.add_series('Missing',missing)
                ch=slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(0.75), Inches(3.15), Inches(11.3), Inches(2.75), cd).chart; ch.has_legend=True; ch.legend.position=XL_LEGEND_POSITION.BOTTOM
                for i,c in enumerate([BLUE,GOLD]): ch.series[i].format.fill.solid(); ch.series[i].format.fill.fore_color.rgb=rgb(c)
        elif typ=='actions':
            rows=s.get('data',[])
            y=1.55
            for r in rows[:6]:
                box(slide,0.75,y,11.75,0.72,WHITE,MID)
                text(slide,0.95,y+0.11,0.75,0.25,str(r.get('priority','Action')),8,True,GOLD)
                text(slide,1.7,y+0.08,9.95,0.48,str(r.get('recommendation','')),10,False,DARK)
                y+=0.82
            if not rows: text(slide,0.8,1.65,11.5,0.6,'No evidence-backed recommendations were produced from the available data.',11,False,'687783')
        elif typ=='quality':
            rows=s.get('data',[])
            y=1.55
            if not rows: text(slide,0.8,y,11.5,0.6,'No mapped operational quality findings were detected.',11,False,GREEN)
            for r in rows[:7]:
                box(slide,0.75,y,11.75,0.68,WHITE,MID)
                text(slide,0.95,y+0.09,1.0,0.25,str(r.get('severity','Review')),8,True,RED if str(r.get('severity','')).lower()=='high' else GOLD)
                text(slide,1.9,y+0.08,10.25,0.48,str(r.get('finding',r.get('rule',''))),9,False,DARK); y+=0.77
        elif typ=='method':
            box(slide,0.75,1.55,11.75,4.8,LIGHT,LIGHT)
            text(slide,1.0,1.85,10.9,0.35,'How to verify this report',13,True,NAVY)
            y=2.38
            for item in s.get('items',[]):
                text(slide,1.0,y,0.3,0.25,'•',10,True,GOLD); text(slide,1.3,y,10.55,0.48,str(item).capitalize(),9,False,DARK); y+=0.48
            text(slide,1.0,5.75,10.8,0.35,'Detailed calculations, source registers, exceptions and traceability belong in the analytical addendum.',9,False,'5C6B78')
        elif typ=='voc':
            data=s.get('data') or {}
            box(slide,0.8,1.6,3.1,2.0,WHITE,MID); text(slide,1.05,1.95,2.5,0.25,'VoC responses',8,False,'687783'); text(slide,1.05,2.32,2.5,0.55,str(data.get('responses',0)),24,True,BLUE)
            box(slide,4.15,1.6,3.1,2.0,WHITE,MID); text(slide,4.4,1.95,2.5,0.25,'Positive / resolved',8,False,'687783'); text(slide,4.4,2.32,2.5,0.55,f"{data.get('positive_pct','Not available')}%",24,True,GREEN)
            text(slide,0.85,4.15,11.2,0.7,'VoC is shown only when supplied evidence is available and labelled for the selected reporting context.',10,False,DARK)
        footer(slide,dept,team,period)
    prs.save(out_path)
