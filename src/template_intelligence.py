from __future__ import annotations
from pathlib import Path
import re
from typing import Any
from pptx import Presentation

SECTION_RULES = [
    ("Executive Coverage Snapshot", ["valid cases", "wards", "coverage", "missing wards", "movement"]),
    ("Ward Coverage by Municipality / Region", ["ward", "municipality", "region", "target", "covered", "missing"]),
    ("Channel and Time Coverage", ["channel", "hour", "time", "cases", "new wards"]),
    ("Services and Demand Mix", ["service", "demand", "category", "cases"]),
    ("Voice of the Citizen and Resolution Position", ["voc", "responses", "happiness", "resolved", "survey", "resolution"]),
    ("Actions, Data Quality and Management Decisions", ["action", "owner", "data quality", "decision", "caveat", "addendum"]),
]

KEYWORD_MAP = {
    "coverage": ["ward", "coverage", "municipality", "region"],
    "channel": ["channel"],
    "time": ["time", "hour", "date", "period"],
    "services": ["service", "demand", "category"],
    "voc": ["voc", "voice of citizen", "happiness", "satisfaction", "survey", "nps"],
    "resolution": ["resolution", "resolved", "status", "case result"],
    "data quality": ["data quality", "missing", "duplicate", "validation", "correction"],
    "actions": ["action", "recommendation", "decision", "owner", "escalation"],
}


def extract_template_sections(template_path: str | Path) -> list[dict[str, Any]]:
    prs = Presentation(str(template_path))
    sections = []
    for idx, slide in enumerate(prs.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text.strip())
        title = next((line.strip() for line in text.splitlines() if re.match(r"^\d+\.\s+", line.strip())), f"Slide {idx}")
        sections.append({"slide": idx, "title": title, "text": text})
    return sections


def _available_capabilities(analysis: dict[str, Any]) -> set[str]:
    cols = analysis.get("columns", {})
    caps = set()
    if cols.get("ward"): caps.add("coverage")
    if cols.get("channel"): caps.add("channel")
    if analysis.get("dates", {}).get("available"): caps.add("time")
    if any(cols.get(k) for k in ["category1", "category2", "category3", "category4", "category5"]): caps.add("services")
    if cols.get("status"): caps.add("resolution")
    if cols.get("case_result") or cols.get("resolution_comments"): caps.add("resolution")
    if analysis.get("quality") is not None: caps.add("data quality")
    caps.add("actions")
    return caps


def assess_template(template_path: str | Path, analysis: dict[str, Any], requirements: str = "") -> dict[str, Any]:
    sections = extract_template_sections(template_path)
    caps = _available_capabilities(analysis)
    req_lower = requirements.lower()
    rows = []
    for name, keywords in SECTION_RULES:
        matched_slide = next((s["slide"] for s in sections if name.lower() in s["title"].lower()), None)
        relevant = [cap for cap, words in KEYWORD_MAP.items() if any(w in req_lower for w in words)] if requirements else []
        supported = [cap for cap in set(c for c in caps) if any(k in name.lower() for k in ([cap] + keywords))]
        # More explicit section capability mapping
        mapping = {
            "Executive Coverage Snapshot": ["coverage", "time", "actions"],
            "Ward Coverage by Municipality / Region": ["coverage"],
            "Channel and Time Coverage": ["channel", "time"],
            "Services and Demand Mix": ["services"],
            "Voice of the Citizen and Resolution Position": ["voc", "resolution"],
            "Actions, Data Quality and Management Decisions": ["actions", "data quality"],
        }
        needed = mapping.get(name, [])
        available = [x for x in needed if x in caps]
        if len(available) == len(needed) and needed:
            status = "Complete"
        elif available:
            status = "Partial"
        else:
            status = "Gap"
        rows.append({"section": name, "slide": matched_slide, "capabilities_expected": ", ".join(needed), "capabilities_available": ", ".join(available), "coverage": status})

    # Requirement coverage
    req_rows = []
    if requirements.strip():
        bullets = [re.sub(r"^[\-•*\d.)\s]+", "", x).strip() for x in re.split(r"\n|;|(?<=[.!?])\s+", requirements) if x.strip()]
        for i, req in enumerate(bullets, 1):
            terms = [cap for cap, words in KEYWORD_MAP.items() if any(w in req.lower() for w in words)]
            supported = [x for x in terms if x in caps]
            req_rows.append({"id": f"REQ-{i:03d}", "requirement": req, "mapped_capabilities": ", ".join(terms), "data_support": "Supported" if terms and len(supported) == len(terms) else ("Partial" if supported else "Not established")})

    score = {"Complete": 1.0, "Partial": 0.5, "Gap": 0.0}
    section_score = round(100 * sum(score[r["coverage"]] for r in rows) / len(rows), 1) if rows else 0
    recs = []
    for r in rows:
        if r["coverage"] == "Gap":
            recs.append(f"Add or strengthen supporting analysis for the '{r['section']}' section before treating the template as fully covered.")
        elif r["coverage"] == "Partial":
            recs.append(f"Review the '{r['section']}' section: some expected analytical capabilities are not currently supported by the dataset.")
    if requirements and any(r["data_support"] != "Supported" for r in req_rows):
        recs.append("Review user requirements marked Partial or Not established; do not imply that unavailable fields are supported by the dataset.")
    return {"sections": rows, "requirements": req_rows, "score": section_score, "recommendations": recs, "template_slides": sections}
